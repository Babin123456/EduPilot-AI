"""AI routes — EduPilot AI chat, RAG document management, and generation endpoints."""

from __future__ import annotations

import io
import os
import uuid
from datetime import datetime, timedelta, timezone
from pathlib import Path

import httpx
from fastapi import (
    APIRouter,
    BackgroundTasks,
    Depends,
    File,
    HTTPException,
    Request,
    UploadFile,
)
from pydantic import BaseModel
from pymongo.database import Database

from app.api.deps import get_current_teacher
from app.core.config import get_settings
from app.core.database import get_db
from app.core.exceptions import http_400, http_404
from app.models.ai_models import new_ai_conversation, new_ai_message
from app.models.student import student_full_name
from app.services.rag_service import (
    delete_rag_document,
    ingest_document,
    list_rag_documents,
    retrieve_context,
    rewrite_query_with_history,
)

router = APIRouter()


class ChatRequest(BaseModel):
    message: str
    conversation_id: str | None = None
    class_id: str | None = None
    file_context: str | None = None
    image_b64: str | None = None
    mime_type: str | None = None


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# DOCUMENT & FILE PARSER (PPT, Excel, Image — non-RAG)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def parse_uploaded_file(file_bytes: bytes, filename: str, content_type: str) -> dict:
    """Extract text and metadata from uploaded Image, PDF, Excel, or PPT file."""
    ext = os.path.splitext(filename)[1].lower()
    text_content = ""
    file_type = "document"

    try:
        if ext in [".pptx", ".ppt"]:
            file_type = "presentation"
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_bytes))
                slides_data = []
                for i, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    if slide_text:
                        slides_data.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))
                text_content = "\n\n".join(slides_data)
            except Exception:
                text_content = f"PPT Presentation: {filename} ({len(file_bytes)} bytes). Contains slide topics & outlines."

        elif ext in [".xlsx", ".xls", ".csv"]:
            file_type = "spreadsheet"
            if ext == ".csv":
                text_content = file_bytes.decode("utf-8", errors="ignore")[:4000]
            else:
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
                    sheets_summary = []
                    for sheet_name in wb.sheetnames:
                        sheet = wb[sheet_name]
                        rows = list(sheet.iter_rows(values_only=True))
                        if rows:
                            headers = [str(cell) for cell in rows[0] if cell is not None]
                            row_preview = []
                            for row in rows[1:15]:
                                row_str = ", ".join([str(c) for c in row if c is not None])
                                if row_str:
                                    row_preview.append(row_str)
                            sheets_summary.append(
                                f"Sheet: {sheet_name}\nHeaders: {', '.join(headers)}\nSample Rows:\n" + "\n".join(row_preview)
                            )
                    text_content = "\n\n".join(sheets_summary)
                except Exception:
                    text_content = f"Excel Spreadsheet: {filename} ({len(file_bytes)} bytes)."

        elif ext == ".pdf":
            file_type = "pdf"
            try:
                from pypdf import PdfReader
                reader = PdfReader(io.BytesIO(file_bytes))
                pages_text = []
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        pages_text.append(page_text)
                text_content = "\n\n".join(pages_text)[:8000]
            except Exception:
                text_content = f"PDF Document: {filename} ({len(file_bytes)} bytes)."

        elif ext in [".docx", ".doc"]:
            file_type = "docx"
            try:
                import tempfile
                import docx2txt
                with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
                    tmp.write(file_bytes)
                    tmp_path = tmp.name
                try:
                    text_content = docx2txt.process(tmp_path)[:8000]
                finally:
                    if os.path.exists(tmp_path):
                        os.unlink(tmp_path)
            except Exception:
                text_content = f"Word Document: {filename} ({len(file_bytes)} bytes)."

        elif ext in [".png", ".jpg", ".jpeg", ".webp", ".gif"]:
            file_type = "image"
            try:
                import base64
                from PIL import Image
                img = Image.open(io.BytesIO(file_bytes))
                b64_data = base64.b64encode(file_bytes).decode("utf-8")
                mime_type = content_type or f"image/{ext.replace('.', '')}"
                text_content = f"Image File: {filename} ({img.format}, {img.size[0]}x{img.size[1]} px). Base64 encoded for multimodal Gemini 1.5 Flash Vision inspection."
            except Exception:
                text_content = f"Image File: {filename} ({len(file_bytes)} bytes) uploaded and received for visual inspection."

        else:
            text_content = file_bytes.decode("utf-8", errors="ignore")[:4000]

    except Exception:
        text_content = f"File {filename} uploaded successfully ({len(file_bytes)} bytes)."

    return {
        "filename": filename,
        "file_type": file_type,
        "size_bytes": len(file_bytes),
        "text_content": text_content.strip() or f"Content from {filename}",
        "image_b64": b64_data if ext in [".png", ".jpg", ".jpeg", ".webp", ".gif"] and 'b64_data' in locals() else None,
        "mime_type": content_type or f"image/{ext.replace('.', '')}" if ext in [".png", ".jpg", ".jpeg", ".webp", ".gif"] else None,
    }


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# FAST NON-BLOCKING LLM CALLS (Groq / Gemini)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _call_groq_llm(messages: list[dict], api_key: str, model: str) -> str | None:
    """Call Groq API directly using httpx with 5.0s timeout."""
    api_key = (api_key or "").strip()
    if not api_key or "your_" in api_key or len(api_key) < 25:
        return None
    try:
        url = "https://api.groq.com/openai/v1/chat/completions"
        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        }
        payload = {
            "model": model or "llama-3.3-70b-versatile",
            "messages": messages,
            "temperature": 0.7,
            "max_tokens": 1024,
        }
        with httpx.Client(timeout=5.0) as client:
            response = client.post(url, json=payload, headers=headers)
            if response.status_code == 200:
                data = response.json()
                return data["choices"][0]["message"]["content"]
            else:
                print(f"[LLM Error] Groq API returned status {response.status_code}: {response.text}")
    except Exception as exc:
        print(f"[LLM Exception] Groq API call failed: {exc}")
    return None


def _call_gemini_llm(prompt: str, api_key: str, model: str, image_b64: str | None = None, mime_type: str = "image/png") -> str | None:
    """Call Gemini API directly using httpx with inline base64 image data for visual analysis."""
    api_key = (api_key or "").strip()
    if not api_key or "your_" in api_key or len(api_key) < 15:
        print(f"[Gemini] SKIP — key invalid (length={len(api_key)}, starts='{api_key[:6]}...')")
        return None

    # Clean base64 string (strip data URI prefix if present and strip whitespace)
    clean_b64 = None
    if image_b64:
        clean_b64 = image_b64.strip()
        if "," in clean_b64:
            clean_b64 = clean_b64.split(",", 1)[1].strip()

    actual_b64 = clean_b64
    actual_mime = mime_type or "image/jpeg"

    if clean_b64:
        try:
            import base64 as b64mod
            from PIL import Image as PILImage
            raw_bytes = b64mod.b64decode(clean_b64)
            img = PILImage.open(io.BytesIO(raw_bytes))
            # Resize if image is larger than 1024px on any side
            max_dim = 1024
            if img.width > max_dim or img.height > max_dim:
                img.thumbnail((max_dim, max_dim), PILImage.LANCZOS)
            # Convert to JPEG for smaller payload
            buf = io.BytesIO()
            rgb_img = img.convert("RGB") if img.mode != "RGB" else img
            rgb_img.save(buf, format="JPEG", quality=80)
            actual_b64 = b64mod.b64encode(buf.getvalue()).decode("utf-8")
            actual_mime = "image/jpeg"
            print(f"[Gemini] Image compressed: {len(clean_b64)} → {len(actual_b64)} chars (JPEG {img.width}x{img.height})")
        except Exception as compress_err:
            print(f"[Gemini] Image compression warning, using cleaned base64: {compress_err}")

    # Build candidate (api_version, model_name) pairs to try for this key — gemini-2.5-flash is primary verified model
    candidate_endpoints = [
        ("v1beta", "gemini-2.5-flash"),
        ("v1beta", "gemini-2.5-flash-lite"),
        ("v1beta", "gemini-2.5-pro"),
        ("v1beta", "gemini-flash-latest"),
        ("v1beta", "gemini-pro-latest"),
    ]
    # Deduplicate preserving order
    unique_candidates = []
    for ver, mname in candidate_endpoints:
        if (ver, mname) not in unique_candidates:
            unique_candidates.append((ver, mname))

    headers = {"Content-Type": "application/json"}
    parts: list[dict] = []
    if actual_b64:
        clean_mime = actual_mime.split(";")[0].strip() if actual_mime else "image/jpeg"
        parts.append({
            "inline_data": {
                "mime_type": clean_mime,
                "data": actual_b64
            }
        })
    parts.append({"text": prompt})

    payload = {
        "contents": [{"parts": parts}],
        "generationConfig": {
            "temperature": 0.7,
            "maxOutputTokens": 2048,
        },
    }

    b64_len = len(actual_b64) if actual_b64 else 0

    for api_ver, target_model in unique_candidates:
        url = f"https://generativelanguage.googleapis.com/{api_ver}/models/{target_model}:generateContent?key={api_key}"
        print(f"[Gemini] Calling {api_ver}/{target_model} | key=...{api_key[-6:]} | image_b64_len={b64_len} | prompt_len={len(prompt)}")

        try:
            with httpx.Client(timeout=6.0) as client:
                response = client.post(url, json=payload, headers=headers)
                print(f"[Gemini] {api_ver}/{target_model} status={response.status_code}")
                if response.status_code == 200:
                    data = response.json()
                    try:
                        text = data["candidates"][0]["content"]["parts"][0]["text"]
                        print(f"[Gemini] SUCCESS ({api_ver}/{target_model}) — response length={len(text)}")
                        return text
                    except (KeyError, IndexError) as parse_err:
                        if "candidates" in data and data["candidates"]:
                            finish_reason = data["candidates"][0].get("finishReason", "")
                            if finish_reason == "SAFETY":
                                print(f"[Gemini] Blocked by safety filters")
                                return "I received your image but the content was flagged by safety filters. Please try with a different image or question."
                        print(f"[Gemini] Parse error on 200 response: {parse_err}")
                elif response.status_code == 429:
                    print(f"[Gemini] 429 Rate Limit / Quota Exceeded on key ...{api_key[-6:]} — switching key immediately")
                    break  # Key is out of quota, stop trying other models on this key
                else:
                    print(f"[Gemini] API error {api_ver}/{target_model} status={response.status_code}: {response.text[:300]}")
        except httpx.TimeoutException:
            print(f"[Gemini] TIMEOUT on {api_ver}/{target_model} (6s limit) — key=...{api_key[-6:]}")
        except Exception as exc:
            print(f"[Gemini] EXCEPTION on {api_ver}/{target_model}: {type(exc).__name__}: {exc}")

    return None


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# DIAGNOSTIC: API KEY HEALTH & LIVE TEST
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

@router.get("/debug-keys")
async def debug_keys():
    """Diagnostic endpoint to verify API key availability on Vercel (never exposes actual key values)."""
    settings = get_settings()
    def key_status(val: str | None) -> dict:
        v = (val or "").strip()
        return {
            "length": len(v),
            "present": bool(v),
            "valid": bool(v and "your_" not in v and len(v) >= 15),
            "preview": f"...{v[-4:]}" if len(v) >= 15 else "(empty/invalid)",
        }

    return {
        "pydantic_settings": {
            "gemini_api_key": key_status(settings.gemini_api_key),
            "gemini_api_key_1": key_status(settings.gemini_api_key_1),
            "gemini_api_key_2": key_status(settings.gemini_api_key_2),
            "groq_api_key_1": key_status(settings.groq_api_key_1),
            "groq_api_key_2": key_status(settings.groq_api_key_2),
            "gemini_model": settings.gemini_model,
        },
        "os_environ": {
            "GEMINI_API_KEY": key_status(os.environ.get("GEMINI_API_KEY")),
            "GEMINI_API_KEY_1": key_status(os.environ.get("GEMINI_API_KEY_1")),
            "GEMINI_API_KEY_2": key_status(os.environ.get("GEMINI_API_KEY_2")),
            "GROQ_API_KEY_1": key_status(os.environ.get("GROQ_API_KEY_1")),
            "GROQ_API_KEY_2": key_status(os.environ.get("GROQ_API_KEY_2")),
            "VERCEL": os.environ.get("VERCEL", "not set"),
        },
    }


@router.get("/test-gemini")
async def test_gemini():
    """Live diagnostic endpoint: Queries Google for available models and performs test calls to Google Gemini API."""
    settings = get_settings()
    keys = [
        ("gemini_api_key", settings.gemini_api_key or os.environ.get("GEMINI_API_KEY", "")),
        ("gemini_api_key_1", settings.gemini_api_key_1 or os.environ.get("GEMINI_API_KEY_1", "")),
        ("gemini_api_key_2", settings.gemini_api_key_2 or os.environ.get("GEMINI_API_KEY_2", "")),
    ]

    results = []
    test_b64 = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="

    endpoints_to_test = [
        ("v1beta", "gemini-1.5-flash-latest"),
        ("v1beta", "gemini-2.5-flash"),
        ("v1beta", "gemini-1.5-pro-latest"),
        ("v1beta", "gemini-2.0-flash-exp"),
    ]

    for name, key_val in keys:
        k_clean = (key_val or "").strip()
        if not k_clean or len(k_clean) < 15 or "your_" in k_clean:
            results.append({"name": name, "status": "skipped", "reason": "empty or placeholder key"})
            continue

        # Fetch list of available models from Google for this key
        available_models = []
        try:
            list_url = f"https://generativelanguage.googleapis.com/v1beta/models?key={k_clean}"
            with httpx.Client(timeout=4.0) as client:
                res_list = client.get(list_url)
                if res_list.status_code == 200:
                    model_objs = res_list.json().get("models", [])
                    available_models = [m.get("name", "").replace("models/", "") for m in model_objs if "generateContent" in m.get("supportedGenerationMethods", [])]
        except Exception as list_exc:
            available_models = [f"Error listing models: {list_exc}"]

        key_results = []
        for api_ver, model_name in endpoints_to_test:
            try:
                url = f"https://generativelanguage.googleapis.com/{api_ver}/models/{model_name}:generateContent?key={k_clean}"
                payload = {
                    "contents": [{
                        "parts": [
                            {"inline_data": {"mime_type": "image/png", "data": test_b64}},
                            {"text": "Describe this test pixel in one word."}
                        ]
                    }]
                }
                with httpx.Client(timeout=4.0) as client:
                    res = client.post(url, json=payload, headers={"Content-Type": "application/json"})
                    if res.status_code == 200:
                        text = res.json()["candidates"][0]["content"]["parts"][0]["text"]
                        key_results.append({"endpoint": f"{api_ver}/{model_name}", "http_status": 200, "success": True, "response": text.strip()})
                    else:
                        key_results.append({"endpoint": f"{api_ver}/{model_name}", "http_status": res.status_code, "success": False, "error": res.text[:200]})
            except Exception as exc:
                key_results.append({"endpoint": f"{api_ver}/{model_name}", "success": False, "error": f"{type(exc).__name__}: {str(exc)}"})

        results.append({
            "name": name,
            "key_preview": f"...{k_clean[-6:]}",
            "available_models_from_google": available_models[:10],
            "trials": key_results
        })

    return {"test_time": datetime.now(timezone.utc).isoformat(), "results": results}


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# RAG DOCUMENT MANAGEMENT ROUTES
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

@router.post("/rag/upload")
async def upload_rag_document(
    file: UploadFile = File(...),
    teacher: dict = Depends(get_current_teacher),
    db: Database = Depends(get_db),
):
    """Upload a PDF or DOCX file for RAG ingestion (parse → chunk → embed → index).

    The document is split into chunks, each embedded using all-MiniLM-L6-v2,
    and stored in MongoDB for vector similarity search during chat.
    """
    settings = get_settings()
    filename = file.filename or "uploaded_file"
    ext = os.path.splitext(filename)[1].lower()

    # Validate file type
    if ext not in (".pdf", ".docx"):
        raise HTTPException(
            status_code=400,
            detail="Only PDF (.pdf) and DOCX (.docx) files are supported for RAG. "
                   "Use the regular file upload for PPT, Excel, and Image files.",
        )

    # Read and validate file size
    file_bytes = await file.read()
    max_size = settings.rag_max_file_size_mb * 1024 * 1024
    if len(file_bytes) > max_size:
        raise HTTPException(
            status_code=400,
            detail=f"File size ({len(file_bytes) / (1024*1024):.1f}MB) exceeds "
                   f"the maximum allowed size of {settings.rag_max_file_size_mb}MB.",
        )

    # Run ingestion pipeline
    try:
        doc_record = ingest_document(
            file_bytes=file_bytes,
            filename=filename,
            teacher_id=teacher["id"],
            db=db,
        )
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to process document: {str(e)}",
        )

    return {
        "success": True,
        "document": {
            "id": doc_record["id"],
            "filename": doc_record["filename"],
            "file_type": doc_record["file_type"],
            "chunk_count": doc_record["chunk_count"],
            "file_size_bytes": doc_record["file_size_bytes"],
            "status": doc_record["status"],
            "created_at": doc_record["created_at"].isoformat()
            if hasattr(doc_record.get("created_at"), "isoformat")
            else doc_record.get("created_at"),
        },
        "message": f"Document '{filename}' indexed successfully with {doc_record['chunk_count']} chunks.",
    }


@router.get("/rag/documents")
def get_rag_documents(
    teacher: dict = Depends(get_current_teacher),
    db: Database = Depends(get_db),
):
    """List all RAG documents uploaded by the current teacher."""
    docs = list_rag_documents(teacher["id"], db)
    return [
        {
            "id": d["id"],
            "filename": d["filename"],
            "file_type": d["file_type"],
            "chunk_count": d.get("chunk_count", 0),
            "file_size_bytes": d.get("file_size_bytes", 0),
            "status": d.get("status", "unknown"),
            "created_at": d["created_at"].isoformat()
            if hasattr(d.get("created_at"), "isoformat")
            else d.get("created_at"),
        }
        for d in docs
    ]


@router.delete("/rag/documents/{document_id}")
def delete_rag_doc(
    document_id: str,
    teacher: dict = Depends(get_current_teacher),
    db: Database = Depends(get_db),
):
    """Delete a RAG document and all its embedding chunks."""
    success = delete_rag_document(document_id, teacher["id"], db)
    if not success:
        raise http_404("Document not found or you don't have permission to delete it.")
    return {"success": True, "message": "Document and all chunks deleted successfully."}


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# NON-RAG FILE UPLOAD (PPT, Excel, Image)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

@router.post("/upload-file")
async def upload_file_for_ai(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    teacher: dict = Depends(get_current_teacher),
    db: Database = Depends(get_db),
):
    """Upload Image, PDF, Excel, or PPT file for AI analysis, storage, and RAG indexing."""
    settings = get_settings()
    file_bytes = await file.read()
    filename = file.filename or "uploaded_file"
    ext = os.path.splitext(filename)[1].lower()
    parsed = parse_uploaded_file(file_bytes, filename, file.content_type or "")

    image_url = None
    # Save file to disk for media access if it is an image
    if ext in [".png", ".jpg", ".jpeg", ".webp", ".gif"]:
        storage_path = settings.storage_path
        safe_filename = f"chat_img_{uuid.uuid4().hex[:8]}_{filename.replace(' ', '_')}"
        file_dest = storage_path / safe_filename
        with open(file_dest, "wb") as f:
            f.write(file_bytes)
        image_url = f"/media/{safe_filename}"

    # Auto-ingest documents into RAG Knowledge Library (PDF & DOCX)
    if ext in [".pdf", ".docx"]:
        background_tasks.add_task(
            ingest_document,
            file_bytes=file_bytes,
            filename=filename,
            teacher_id=teacher["id"],
            db=db,
        )

    return {
        "success": True,
        "filename": parsed["filename"],
        "file_type": parsed["file_type"],
        "size_bytes": parsed["size_bytes"],
        "extracted_text": parsed["text_content"],
        "image_url": image_url,
        "image_b64": parsed.get("image_b64"),
        "mime_type": parsed.get("mime_type"),
        "summary": f"Uploaded {parsed['file_type'].upper()} file '{parsed['filename']}' ({parsed['size_bytes']} bytes) ready for AI analysis & stored in Knowledge Base.",
    }


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# MAIN CHAT ENDPOINT (with RAG + Conversation History)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

@router.post("/chat")
def chat(
    request: Request,
    body: ChatRequest,
    teacher: dict = Depends(get_current_teacher),
    db: Database = Depends(get_db),
):
    """Send a message to EduPilot AI with RAG retrieval and conversation history.

    Flow:
    1. Load last 20 messages from this conversation
    2. Rewrite query using chat history (for better vector search)
    3. Vector search: retrieve top-k relevant document chunks
    4. Build enhanced system prompt with RAG context + class context
    5. Call LLM (Groq → Gemini → Fallback)
    6. Save messages and return
    """
    settings = get_settings()

    # ── Get active class context ──
    class_context_str = ""
    student_summary_str = ""
    if body.class_id:
        tca = db.teacher_course_assignments.find_one({"id": body.class_id})
        if tca:
            course = db.courses.find_one({"id": tca["course_id"]})
            year = db.years.find_one({"id": tca["year_id"]})
            section = db.sections.find_one({"id": tca["section_id"]})

            class_context_str = f"Active Class: {course['name']} ({course['code']}) • {year['label']} Section {section['name']} • Room: {tca.get('room')}"

            students = list(db.students.find({"section_id": tca["section_id"], "is_active": True}))
            at_risk = [s for s in students if s.get("attendance_percentage", 100) < 75]

            student_details_list = [
                f"{student_full_name(s)} (Roll: {s['roll_number']}, Attendance: {s.get('attendance_percentage', 0)}%, Risk: {s.get('risk_level', 'normal')})"
                for s in students
            ]
            student_summary_str = (
                f"Total Students: {len(students)} | At-Risk (<75% attendance): {len(at_risk)}\n"
                f"Full Class Student Roster with Exact Database Attendance Records:\n"
                + "\n".join(student_details_list)
            )

    # ── Retrieve or create conversation ──
    conversation = None
    if body.conversation_id:
        conversation = db.ai_conversations.find_one({
            "id": body.conversation_id, "teacher_id": teacher["id"],
        })

    if not conversation:
        conversation = new_ai_conversation(
            teacher_id=teacher["id"],
            teacher_course_assignment_id=body.class_id,
            title=body.message[:100],
        )
        db.ai_conversations.insert_one(conversation)

    # ── Load conversation history (last 20 messages) ──
    past_messages = list(
        db.ai_messages.find({"conversation_id": conversation["id"]})
        .sort("created_at", 1)
        .limit(20)
    )

    # ── Save current user message ──
    user_msg = new_ai_message(
        conversation_id=conversation["id"],
        role="user",
        content=body.message,
    )
    db.ai_messages.insert_one(user_msg)

    # ── RAG: Rewrite query with conversation history ──
    chat_history_for_rewrite = [
        {"role": m["role"], "content": m["content"]} for m in past_messages
    ]
    standalone_query = rewrite_query_with_history(
        body.message, chat_history_for_rewrite, settings
    )

    # ── RAG: Retrieve relevant document chunks (skip for casual greetings like "hello", "hi") ──
    casual_greetings = {"hello", "hi", "hey", "good morning", "good afternoon", "good evening", "greetings", "hey there", "hello there"}
    is_casual = body.message.strip().lower().rstrip("!.,") in casual_greetings

    rag_context = ""
    if not is_casual:
        rag_context = retrieve_context(standalone_query, teacher["id"], db)

    # ── Build the full user input ──
    full_user_input = body.message
    if body.file_context:
        full_user_input += f"\n\n[Attached File Content for Analysis]:\n{body.file_context}"

    # ── Build System Context Prompt (with RAG context) ──
    rag_section = ""
    if rag_context:
        rag_section = (
            f"\n[Retrieved Document Context — from teacher's uploaded PDF/DOCX files]:\n"
            f"{rag_context}\n"
        )

    # ── Real-Time Temporal Context (Indian Standard Time - IST, UTC+5:30) ──
    ist_tz = timezone(timedelta(hours=5, minutes=30))
    now_ist = datetime.now(ist_tz)
    current_date_str = now_ist.strftime("%B %d, %Y")
    current_time_str = now_ist.strftime("%I:%M %p IST")
    current_day_name = now_ist.strftime("%A")
    today_dow = now_ist.weekday()

    # Load today's upcoming classes for the teacher
    from app.api.v1.routes.timetable import _get_timetable_for_day
    today_schedule_entries = _get_timetable_for_day(db, teacher, today_dow)
    schedule_summary_lines = []
    if today_schedule_entries:
        for entry in today_schedule_entries:
            attendance_tag = " [Attendance Taken]" if entry.get("attendance_taken") else " [Pending Attendance]"
            schedule_summary_lines.append(
                f"- {entry['start_time']} - {entry['end_time']}: {entry['course_code']} ({entry['course_name']}) "
                f"for {entry['year_label']} Sec {entry['section_name']} in Room {entry['room']}{attendance_tag}"
            )
        upcoming_schedule_str = "\n".join(schedule_summary_lines)
    else:
        upcoming_schedule_str = "No classes scheduled for today (or weekend / non-teaching day)."

    system_prompt = (
        f"You are EduPilot AI, the intelligent academic copilot for university faculty.\n"
        f"You are assisting Professor {teacher['full_name']} ({teacher.get('designation', '')}, {teacher.get('specialization', 'CSE')}).\n"
        f"Real-Time Temporal Context (Indian Standard Time):\n"
        f"  - Date: {current_date_str}\n"
        f"  - Day: {current_day_name}\n"
        f"  - Current Time: {current_time_str}\n"
        f"Today's Teaching Schedule & Upcoming Classes:\n{upcoming_schedule_str}\n"
        f"Current Active Class Context: {class_context_str or 'General Academic Workspace'}\n"
        f"Live Database Information:\n{student_summary_str or 'N/A'}\n"
        f"{rag_section}\n"
        f"Instructions:\n"
        f"1. Always use the real-time temporal context and upcoming schedule above to answer questions like 'What day is it?', 'Do I have any classes today?', 'What is my next class?', or 'What time is it?'.\n"
        f"2. Always use the live student database information provided above to give exact student names, roll numbers, and attendance percentages when asked.\n"
        f"3. Remember previous user questions and context in this conversation thread.\n"
        f"4. When answering questions about uploaded documents, cite the specific source filename and page numbers from the Retrieved Document Context above.\n"
        f"5. Format responses cleanly using rich GitHub Markdown: use **bold** for key metrics/names, *italics* for emphasis, standard markdown `[link label](url)` for web URLs/references, bullet lists, and clean Markdown tables for tabular data.\n"
        f"6. STRICT EMOJI POLICY: DO NOT include any text emojis or emoticons (such as 📚, 📄, 🤖, 😀, etc.) in your responses under any circumstances. Keep responses clean, elegant, and professional.\n"
        f"7. TEMPORAL AWARENESS & REAL-WORLD CONTEXT: The current year is {now_ist.year} (2026). Treat all events, movies, releases, and milestones scheduled for or occurring by 2026 (such as the Project Hail Mary film released in March 2026, directed by Phil Lord & Christopher Miller, starring Ryan Gosling) as current or completed events rather than future/unannounced projects."
    )

    # ── Build LLM message array ──
    llm_messages = [{"role": "system", "content": system_prompt}]
    for pm in past_messages:
        llm_messages.append({"role": pm["role"], "content": pm["content"]})
    llm_messages.append({"role": "user", "content": full_user_input})

    # ── Smart Model Routing: Check if prompt contains attached file/image context ──
    has_attachment = bool((body.file_context and body.file_context.strip()) or body.image_b64)
    is_image = bool(body.image_b64 or "Image File:" in (body.file_context or ""))

    model_used = ""
    ai_response = None

    # Helper for multi-key Gemini failover
    def _gemini_with_failover(prompt_text: str, b64_img: str | None = None, mime: str | None = None) -> str | None:
        # Collect keys from both pydantic-settings and direct os.environ (Vercel injects into os.environ)
        keys_to_try = [
            ("settings.gemini_api_key", settings.gemini_api_key),
            ("settings.gemini_api_key_1", settings.gemini_api_key_1),
            ("settings.gemini_api_key_2", settings.gemini_api_key_2),
            ("os.environ GEMINI_API_KEY", os.environ.get("GEMINI_API_KEY")),
            ("os.environ GEMINI_API_KEY_1", os.environ.get("GEMINI_API_KEY_1")),
            ("os.environ GEMINI_API_KEY_2", os.environ.get("GEMINI_API_KEY_2")),
        ]

        # Debug: log which key sources have values
        for source, val in keys_to_try:
            v = (val or "").strip()
            print(f"[Gemini Failover] {source}: len={len(v)}, valid={bool(v and 'your_' not in v and len(v) >= 15)}")

        # Filter non-empty unique keys
        valid_keys = []
        for _source, k in keys_to_try:
            k_clean = (k or "").strip()
            if k_clean and "your_" not in k_clean and len(k_clean) >= 15 and k_clean not in valid_keys:
                valid_keys.append(k_clean)

        print(f"[Gemini Failover] Total valid unique keys: {len(valid_keys)} | has_image={bool(b64_img)} | prompt_len={len(prompt_text)}")

        if not valid_keys:
            print("[Gemini Failover] NO VALID KEYS FOUND — falling back to offline response generator")
            return None

        for i, key in enumerate(valid_keys):
            print(f"[Gemini Failover] Trying key {i+1}/{len(valid_keys)} (ends ...{key[-6:]})")
            res = _call_gemini_llm(
                prompt=prompt_text,
                api_key=key,
                model=settings.gemini_model or "gemini-2.0-flash-lite",
                image_b64=b64_img,
                mime_type=mime or "image/jpeg",
            )
            if res:
                print(f"[Gemini Failover] Key {i+1} succeeded!")
                return res
            print(f"[Gemini Failover] Key {i+1} failed, trying next...")
        print("[Gemini Failover] ALL keys exhausted — no response obtained")
        return None

    if has_attachment:
        # Attachment detected (Image or Document) -> Smart Route to Gemini Vision / Flash first with Multi-Key Failover
        model_used = "Gemini 1.5 Flash (Vision & Attachment Route)"
        if body.image_b64:
            gemini_prompt = (
                f"You are EduPilot AI, an expert visual & document AI assistant for Professor {teacher['full_name']}.\n"
                f"The user has attached an image file for visual inspection.\n"
                f"Inspect the image pixels thoroughly. Read all text, logos, headers, links, table data, and details inside the image.\n"
                f"User Question: {body.message}\n"
                f"Format your response in clean, professional Markdown."
            )
        else:
            gemini_prompt = f"{system_prompt}\n\nUser Question: {full_user_input}"
        ai_response = _gemini_with_failover(gemini_prompt, body.image_b64, body.mime_type)

        # Fallback to Groq ONLY if it's a document (not a visual image requiring Gemini Vision)
        if not ai_response and not body.image_b64:
            model_used = "Groq Llama-3.3-70B (Attachment Fallback)"
            ai_response = _call_groq_llm(llm_messages, settings.groq_api_key_1, settings.groq_model)
            if not ai_response:
                ai_response = _call_groq_llm(llm_messages, settings.groq_api_key_2, settings.groq_model)
    else:
        # Standard Chat -> Primary Route to Groq for ultra-fast response
        model_used = "Groq Llama-3.3-70B"
        ai_response = _call_groq_llm(llm_messages, settings.groq_api_key_1, settings.groq_model)
        if not ai_response:
            ai_response = _call_groq_llm(llm_messages, settings.groq_api_key_2, settings.groq_model)

        # Fallback to Gemini with Multi-Key Failover if Groq fails
        if not ai_response:
            model_used = "Gemini 1.5 Flash (Fallback)"
            gemini_prompt = f"{system_prompt}\n\nUser Question: {full_user_input}"
            ai_response = _gemini_with_failover(gemini_prompt, None, None)

    # ── Fallback to Smart Contextual Response Generator ──
    if not ai_response:
        model_used = "EduPilot Smart Engine"
        ai_response = _generate_contextual_response(body.message, teacher, db, body.class_id, body.file_context, rag_context)

    # ── Determine if RAG was used ──
    content_type = "rag" if rag_context else "text"

    # ── Save assistant message ──
    assistant_msg = new_ai_message(
        conversation_id=conversation["id"],
        role="assistant",
        content=ai_response,
        model_used=model_used,
        content_type=content_type,
    )
    db.ai_messages.insert_one(assistant_msg)

    db.ai_conversations.update_one(
        {"id": conversation["id"]},
        {"$inc": {"message_count": 2}, "$set": {"updated_at": datetime.now(timezone.utc)}},
    )

    # ── Build source references for the frontend ──
    sources = []
    if rag_context:
        # Extract unique source filenames from the context
        import re as re_mod
        source_matches = re_mod.findall(r"from '([^']+)'", rag_context)
        seen = set()
        for src in source_matches:
            if src not in seen:
                sources.append(src)
                seen.add(src)

    return {
        "conversation_id": conversation["id"],
        "message": {
            "id": assistant_msg["id"],
            "role": "assistant",
            "content": ai_response,
            "content_type": content_type,
            "model_used": model_used,
            "sources": sources,
            "created_at": assistant_msg["created_at"].isoformat() if hasattr(assistant_msg.get("created_at"), 'isoformat') else assistant_msg.get("created_at"),
        },
    }


def _generate_contextual_response(
    message: str, teacher: dict, db: Database, class_id: str | None, file_context: str | None, rag_context: str | None = None
) -> str:
    """Smart contextual answer engine when external LLM API keys are unconfigured."""
    msg_lower = message.lower()

    # If RAG context is available, present it
    if rag_context:
        return (
            f"**Document Analysis (from your uploaded Knowledge Base):**\n\n"
            f"Based on your uploaded documents, here is the relevant information:\n\n"
            f"{rag_context[:2000]}\n\n"
            f"---\n"
            f"*Note: For more detailed analysis, configure your Groq or Gemini API keys "
            f"to enable the full AI engine.*"
        )

    if file_context:
        is_image = "Image File:" in file_context or any(kw in file_context.lower() for kw in ["png", "jpg", "jpeg", "webp", "gif"])
        if is_image:
            return (
                f"### Image Received ({file_context.splitlines()[0] if file_context.splitlines() else 'Uploaded Image'})\n\n"
                f"EduPilot AI received your image file and prepared the visual payload, "
                f"but the Gemini Vision API could not process it at this time.\n\n"
                f"**Possible causes:**\n"
                f"- Gemini API keys may have exhausted their free quota (check [Google AI Studio](https://aistudio.google.com/apikey))\n"
                f"- The image may be too large or the API request timed out\n"
                f"- The API key may not have the Gemini Vision model enabled\n\n"
                f"**Try:** Upload a smaller image, or check your API key quota in Google AI Studio."
            )
        else:
            return (
                f"### Document Analysis & Extracted Text\n\n"
                f"{file_context[:3000]}\n\n"
                f"---\n"
                f"**Key Insights & Educational Relevance:**\n"
                f"- Document parsed successfully into readable text context.\n"
                f"- You can ask direct questions about this document or dispatch it to students in **Communications** or **Document Studio**."
            )

    if msg_lower in {"hello", "hi", "hey", "good morning", "good afternoon", "good evening", "greetings", "hey there", "hello there"}:
        return (
            f"Hello Professor {teacher.get('last_name', '')}! Good to see you.\n\n"
            f"How can I assist you with your active class, student attendance, or academic materials today?"
        )

    if "attendance" in msg_lower and any(kw in msg_lower for kw in ["below", "risk", "<", "75", "less", "shortage"]):
        students = []
        # Fallback to teacher's first assigned class if class_id is not passed
        target_tca = None
        if class_id:
            target_tca = db.teacher_course_assignments.find_one({"id": class_id})
        if not target_tca:
            target_tca = db.teacher_course_assignments.find_one({"teacher_id": teacher["id"]})

        if target_tca:
            all_section_students = list(
                db.students.find({
                    "section_id": target_tca["section_id"],
                    "is_active": True,
                })
            )
            # Strict comparison: attendance_percentage strictly less than 75
            students = [s for s in all_section_students if s.get("attendance_percentage", 100) < 75.0]
            students.sort(key=lambda x: x.get("attendance_percentage", 0))

        if students:
            lines = [f"Based on the live student database, here are the **{len(students)} students** in your active class section with attendance below **75%**:\n"]
            for i, s in enumerate(students, 1):
                att = s.get('attendance_percentage', 0)
                lines.append(f"{i}. **{student_full_name(s)}** (Roll: `{s['roll_number']}`, Attendance: **{att:.1f}%**) — Email: `{s['email']}`")
            lines.append(f"\n**Total At-Risk:** {len(students)} students require attention.")
            lines.append("*Tip: You can navigate to the **Communications** module to send automated warning emails to these students.*")
            return "\n".join(lines)
        return "Great news! All students in the active class section have attendance at or above the 75% threshold."

    if any(kw in msg_lower for kw in ["schedule", "routine", "timetable", "today", "classes"]):
        tca_ids = [t["id"] for t in db.teacher_course_assignments.find({"teacher_id": teacher["id"]})]
        entries = list(
            db.timetable_entries.find({
                "teacher_course_assignment_id": {"$in": tca_ids},
            }).sort("start_time", 1)
        )
        if entries:
            lines = ["**Your Academic Teaching Schedule:**\n"]
            days = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday"]
            for e in entries[:6]:
                tca = db.teacher_course_assignments.find_one({"id": e["teacher_course_assignment_id"]})
                course = db.courses.find_one({"id": tca["course_id"]}) if tca else None
                year = db.years.find_one({"id": tca["year_id"]}) if tca else None
                sec = db.sections.find_one({"id": tca["section_id"]}) if tca else None
                day_name = days[e.get("day_of_week", 0)]
                lines.append(f"- **{day_name} {e['start_time']} - {e['end_time']}**: {course['name'] if course else 'Course'} (`{course['code'] if course else ''}`) — {year['label'] if year else ''} Sec {sec['name'] if sec else ''} (Room {e.get('room', '101')})")
            lines.append("\nView full weekly grid under the **Timetable** section.")
            return "\n".join(lines)
        return "You have no classes scheduled for today. Check the **Timetable** module for full weekly routine."

    if "tcp" in msg_lower or "osi" in msg_lower or ("protocol" in msg_lower and "suite" in msg_lower):
        return (
            "### Comprehensive Comparison: TCP/IP Protocol Suite vs. OSI Reference Model\n\n"
            "Both models define standards for network communication, but they differ in structure and practical implementation:\n\n"
            "| Feature | OSI Model (Open Systems Interconnection) | TCP/IP Model (Transmission Control Protocol/Internet Protocol) |\n"
            "| :--- | :--- | :--- |\n"
            "| **Layers** | **7 Layers** (Application, Presentation, Session, Transport, Network, Data Link, Physical) | **4 Layers** (Application, Transport, Internet, Network Access) |\n"
            "| **Approach** | Conceptual & Theoretical reference model | Implementation-oriented, practical Internet architecture |\n"
            "| **Protocols** | Protocol-independent (developed before protocols) | Protocol-dependent (built around TCP and IP) |\n"
            "| **Session & Presentation** | Separate dedicated layers (Layers 5 & 6) | Combined directly into the Application layer |\n"
            "| **Transport Services** | Supports both Connection-Oriented and Connectionless | Supports both (TCP for connection-oriented, UDP for connectionless) |\n\n"
            "#### Key Takeaway for Computer Networks Lecture:\n"
            "- **OSI** is ideal for teaching architectural layering and modular network design.\n"
            "- **TCP/IP** is the actual operational suite powering global Internet traffic today."
        )

    if "quiz" in msg_lower or "outline" in msg_lower or "operating system" in msg_lower:
        return (
            "### Operating Systems (OS) Quiz Topic Outline\n\n"
            "Here is a 4-module quiz outline for your **Operating Systems** course:\n\n"
            "#### Module 1: Process Management & CPU Scheduling\n"
            "- Process states, PCB (Process Control Block), and context switching.\n"
            "- Preemptive vs Non-preemptive scheduling algorithms (FCFS, SJF, Round Robin, Priority).\n"
            "- Multithreading models and race condition prevention.\n\n"
            "#### Module 2: Process Synchronization & Deadlocks\n"
            "- Critical section problem, Peterson's solution, Mutexes, and Semaphores.\n"
            "- Coffman conditions for deadlocks.\n"
            "- Banker's Algorithm for deadlock avoidance and resource allocation graphs.\n\n"
            "#### Module 3: Memory Management & Virtual Memory\n"
            "- Contiguous allocation, Paging, Segmentation, and TLB (Translation Lookaside Buffer).\n"
            "- Page replacement algorithms (FIFO, LRU, Optimal) and Thrashing.\n\n"
            "#### Module 4: File Systems & Storage Management\n"
            "- Directory structures, inode allocation, and disk scheduling algorithms (SSTF, SCAN, C-SCAN).\n\n"
            "*Tip: You can use the **Document Studio** or **Assignments** module to generate downloadable PDF/Word question papers from this outline.*"
        )

    if any(kw in msg_lower for kw in ["what is", "explain", "how does", "difference", "algorithm", "network", "operating system", "data structure", "protocol"]):
        return (
            f"**Academic Topic Breakdown: {message.title()}**\n\n"
            f"Here is a structured breakdown of **{message}**:\n\n"
            f"### 1. Core Concept & Definition\n"
            f"{message} is a foundational subject concept in Computer Science & Engineering, governing system execution, algorithmic logic, and operational performance.\n\n"
            f"### 2. Key Technical Principles\n"
            f"- **Efficiency**: Optimizes processing cycles and throughput.\n"
            f"- **Scalability**: Handles expanding workloads and data flow.\n"
            f"- **Reliability**: Guarantees deterministic outcome and error control.\n\n"
            f"### 3. Practical Industry Applications\n"
            f"- Enterprise software architecture and backend systems\n"
            f"- Distributed cloud computing infrastructure\n"
            f"- Real-time data processing and analytics\n\n"
            f"*Recommendation: Use **Daily Notes** to generate discussion summaries or **Document Studio** to export a PDF quiz on this topic.*"
        )

    return (
        f"Welcome Professor {teacher.get('last_name', '')}. I am **EduPilot AI**, your institutional copilot.\n\n"
        f"### How I Can Assist You:\n"
        f"- **Class Metrics**: Ask *'Which students in my active class have attendance below 75%?'*\n"
        f"- **Schedule & Routine**: Ask *'What classes do I have scheduled for today?'*\n"
        f"- **Document Analysis (RAG)**: Upload any PDF or DOCX using the 📎 button — I'll index it and answer questions about the content.\n"
        f"- **File Analysis**: Upload PPT, Excel, or Image files for instant extraction.\n"
        f"- **Academic Material**: Ask *'Explain the TCP/IP protocol suite vs OSI model'* or *'Generate a quiz topic outline for Operating Systems'*.\n"
        f"- **Daily Notes & Publish**: Use **Daily Notes** or **Document Studio** to generate and dispatch materials."
    )



@router.get("/conversations")
def list_conversations(
    teacher: dict = Depends(get_current_teacher),
    db: Database = Depends(get_db),
):
    """List teacher's AI conversations."""
    convos = list(
        db.ai_conversations.find({"teacher_id": teacher["id"]})
        .sort("updated_at", -1)
        .limit(20)
    )
    return [
        {
            "id": c["id"],
            "title": c.get("title", "New Conversation"),
            "message_count": c.get("message_count", 0),
            "created_at": c["created_at"].isoformat() if hasattr(c.get("created_at"), 'isoformat') else c.get("created_at"),
            "updated_at": c["updated_at"].isoformat() if hasattr(c.get("updated_at"), 'isoformat') else c.get("updated_at"),
        }
        for c in convos
    ]


@router.get("/conversations/{conversation_id}")
def get_conversation(
    conversation_id: str,
    teacher: dict = Depends(get_current_teacher),
    db: Database = Depends(get_db),
):
    """Get full conversation with messages."""
    convo = db.ai_conversations.find_one({"id": conversation_id})
    if not convo:
        raise http_404("Conversation not found")
    if convo["teacher_id"] != teacher["id"]:
        raise http_400("Not authorized")

    messages = list(
        db.ai_messages.find({"conversation_id": conversation_id})
        .sort("created_at", 1)
    )

    return {
        "id": convo["id"],
        "title": convo.get("title", "New Conversation"),
        "messages": [
            {
                "id": m["id"],
                "role": m["role"],
                "content": m["content"],
                "content_type": m.get("content_type", "text"),
                "model_used": m.get("model_used"),
                "created_at": m["created_at"].isoformat() if hasattr(m.get("created_at"), 'isoformat') else m.get("created_at"),
            }
            for m in messages
        ],
    }


@router.delete("/conversations/{conversation_id}")
def delete_conversation(
    conversation_id: str,
    teacher: dict = Depends(get_current_teacher),
    db: Database = Depends(get_db),
):
    """Delete a conversation and its messages."""
    convo = db.ai_conversations.find_one({
        "id": conversation_id, "teacher_id": teacher["id"],
    })
    if not convo:
        raise http_404("Conversation not found")

    db.ai_messages.delete_many({"conversation_id": conversation_id})
    db.ai_conversations.delete_one({"id": conversation_id})
    return {"success": True, "message": "Conversation deleted successfully"}
